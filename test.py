from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------
# Конфигурация
# ---------------------
LOGO_MIREA = "logo_mirea.png"
LOGO_ALTAIR = "logo_altair.png"
TEAM_PHOTO = "team_photo.png"
SCREENSHOTS = ["main_page.png", "profile_page.png"]
COMPETITORS = [
    "Жалялетдинова Динара",
    "Кузнецов Никита", 
    "Прохоров Иван",
    "Тапилин Виктор"
]

# ---------------------
# Основные функции
# ---------------------
def add_slide_title(prs):
    """Титульный слайд"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "UltimateUnity\nПлатформа для IT-проектов"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    # Логотипы
    slide.shapes.add_picture(LOGO_MIREA, Cm(2), Cm(16), height=Cm(2))
    slide.shapes.add_picture(LOGO_ALTAIR, Cm(20), Cm(16), height=Cm(2))

def add_competitors_slide(prs):
    """Слайд с конкурентами"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Конкуренты"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Основные конкуренты на рынке:"
    p.font.bold = True
    
    for competitor in COMPETITORS:
        p = tf.add_paragraph()
        p.text = f"• {competitor}"
        p.level = 1

def add_technology_stack(prs):
    """Стек технологий"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Технологии"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.text = "Backend:\n• Flask\n• SQLAlchemy\n\nFrontend:\n• HTML/CSS\n• JavaScript"

def add_team_slide(prs):
    """Команда разработки"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Наша команда"
    
    # Текст
    textbox = slide.shapes.add_textbox(Cm(2), Cm(4), Cm(12), Cm(10))
    tf = textbox.text_frame
    members = [
        "Емельянов Иван - Team Lead",
        "Буланов Василий - FullStack",
        "Ласиков Богдан - QA Engineer"
    ]
    for member in members:
        p = tf.add_paragraph()
        p.text = member
        p.font.size = Pt(18)
    
    # Фото
    try:
        slide.shapes.add_picture(TEAM_PHOTO, Cm(16), Cm(4), width=Cm(12))
    except FileNotFoundError:
        print("Фото команды не найдено!")

def add_screenshots_slide(prs):
    """Скриншоты интерфейса"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Интерфейс платформы"
    
    for idx, img in enumerate(SCREENSHOTS):
        try:
            slide.shapes.add_picture(img, Cm(2 + idx*15), Cm(4), width=Cm(12))
        except FileNotFoundError:
            print(f"Скриншот {img} не найден")

# ---------------------
# Генерация презентации
# ---------------------
def create_presentation():
    prs = Presentation()
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)
    
    # Добавляем слайды
    add_slide_title(prs)
    add_competitors_slide(prs)
    add_technology_stack(prs)
    add_team_slide(prs)
    add_screenshots_slide(prs)
    
    # Сохраняем
    prs.save("UltimateUnity_Final.pptx")
    print("Презентация создана!")

if __name__ == "__main__":
    create_presentation()